"""
Build thesis defense presentation from template.
Keeps all styles/backgrounds from the template, only changes text and adds images.
"""
import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from lxml import etree

TEMPLATE = r"C:\Users\lividcoffe\Desktop\Шаблон презентации кафедра ВС.pptx"
OUTPUT = r"C:\Users\lividcoffe\Desktop\Презентация_Бесчетнов_ИС241.pptx"
IMG_DIR = r"C:\Users\lividcoffe\Desktop\study\rule-based-UI-generation\slide_images"

prs = Presentation(TEMPLATE)

# --- Helper functions ---

def clone_slide(prs, template_index):
    """Clone a slide from the presentation by index."""
    template_slide = prs.slides[template_index]
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Remove default placeholders content from new slide
    for ph in new_slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)

    # Copy all shapes from template
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # Copy relationships (images etc)
    for rel in template_slide.part.rels.values():
        if "image" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)

    return new_slide


def duplicate_content_slide(prs, source_index=1):
    """Duplicate slide at source_index (0-based). Returns the new slide."""
    return clone_slide(prs, source_index)


def set_title(slide, text):
    """Set the title text of a slide, preserving formatting."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            ph = shape.placeholder_format
            if ph is not None:
                ph_type = ph.type
                if ph_type in (1, 15):  # TITLE, CENTER_TITLE
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        runs = list(para.runs)
                        if runs:
                            runs[0].text = text
                            for extra_run in runs[1:]:
                                extra_run._r.getparent().remove(extra_run._r)
                            return
                    if tf.paragraphs:
                        para = tf.paragraphs[0]
                        run = para.add_run()
                        run.text = text
                    return
        except (ValueError, AttributeError):
            pass
    # Fallback: find first text shape that looks like a title (large font)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rpr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                if rpr is not None:
                    sz = rpr.get('sz')
                    if sz and int(sz) >= 3000:
                        runs = list(para.runs)
                        runs[0].text = text
                        for extra_run in runs[1:]:
                            extra_run._r.getparent().remove(extra_run._r)
                        return


def set_body_text(slide, lines, is_bullet=True):
    """Set body text. lines is a list of strings."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            ph = shape.placeholder_format
        except (ValueError, AttributeError):
            continue
        if ph is None or ph.type != 2:
            continue
        tf = shape.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            template_para = copy.deepcopy(tf.paragraphs[0]._p)
        else:
            template_para = None
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        for p in list(tf._txBody.findall(f'{ns}p')):
            tf._txBody.remove(p)
        for line in lines:
            if template_para is not None:
                new_p = copy.deepcopy(template_para)
                for r in new_p.findall(f'{ns}r'):
                    t = r.find(f'{ns}t')
                    if t is not None:
                        t.text = line
                        break
                else:
                    r_elem = etree.SubElement(new_p, f'{ns}r')
                    rpr = etree.SubElement(r_elem, f'{ns}rPr')
                    rpr.set('lang', 'ru-RU')
                    rpr.set('sz', '2000')
                    t = etree.SubElement(r_elem, f'{ns}t')
                    t.text = line
                tf._txBody.append(new_p)
            else:
                p = tf._txBody.makeelement(f'{ns}p', {})
                r = etree.SubElement(p, f'{ns}r')
                rpr = etree.SubElement(r, f'{ns}rPr')
                rpr.set('lang', 'ru-RU')
                rpr.set('sz', '2000')
                t = etree.SubElement(r, f'{ns}t')
                t.text = line
                tf._txBody.append(p)
        return


def set_page_number(slide, num):
    """Set page number in the small text box."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        is_placeholder = False
        try:
            ph = shape.placeholder_format
            if ph is not None:
                is_placeholder = True
        except (ValueError, AttributeError):
            pass
        if is_placeholder:
            continue
        text = shape.text_frame.text.strip()
        if text.isdigit() and len(text) <= 3:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = str(num)
                    return


def resize_body(slide, new_width=None, new_height=None):
    """Resize the body placeholder text box."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            ph = shape.placeholder_format
            if ph is not None and ph.type == 2:
                if new_width:
                    shape.width = new_width
                if new_height:
                    shape.height = new_height
                return
        except (ValueError, AttributeError):
            pass


def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Add an image to a slide."""
    if width and height:
        slide.shapes.add_picture(image_path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(image_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        slide.shapes.add_picture(image_path, left, top)


# =========================================================
# Build presentation
# =========================================================

# First, let's understand the template structure
# Slide 0 (index 0): Title slide
# Slide 1 (index 1): "Цель работы" - content template
# Slide 2 (index 2): "Постановка задачи" - content template
# Slide 3 (index 3): "Описание разработки" - content template
# Slide 4 (index 4): "Заключение" - content template
# Slide 5 (index 5): "Спасибо за внимание" - final slide

# Strategy: Edit existing slides in place, and add new slides by duplicating

# ---- SLIDE 1: Title ----
slide1 = prs.slides[0]
# Update title
for shape in slide1.shapes:
    if shape.has_text_frame:
        full_text = shape.text_frame.text
        if "Название" in full_text and "квалификационной" in full_text:
            # This is the main title
            ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            txBody = shape.text_frame._txBody
            # Remove all paragraphs
            for p in list(txBody.findall(f'{ns}p')):
                txBody.remove(p)
            # Add new title
            p = etree.SubElement(txBody, f'{ns}p')
            pPr = etree.SubElement(p, f'{ns}pPr')
            pPr.set('algn', 'l')
            lnSpc = etree.SubElement(pPr, f'{ns}lnSpc')
            spcPct = etree.SubElement(lnSpc, f'{ns}spcPct')
            spcPct.set('val', '90000')
            buNone = etree.SubElement(pPr, f'{ns}buNone')

            r = etree.SubElement(p, f'{ns}r')
            rPr = etree.SubElement(r, f'{ns}rPr')
            rPr.set('lang', 'ru-RU')
            latin = etree.SubElement(rPr, f'{ns}latin')
            latin.set('typeface', 'Montserrat')
            ea = etree.SubElement(rPr, f'{ns}ea')
            ea.set('typeface', 'Montserrat')
            cs = etree.SubElement(rPr, f'{ns}cs')
            cs.set('typeface', 'Montserrat')
            t = etree.SubElement(r, f'{ns}t')
            t.text = "Разработка приложения для генерации клиентского интерфейса"

        elif "Выполнил" in full_text:
            # Author/supervisor block
            ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            txBody = shape.text_frame._txBody
            for p in list(txBody.findall(f'{ns}p')):
                txBody.remove(p)

            def add_author_line(txBody, text, bold=False, lang='ru-RU'):
                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                p = etree.SubElement(txBody, f'{ns}p')
                pPr = etree.SubElement(p, f'{ns}pPr')
                pPr.set('algn', 'l')
                lnSpc = etree.SubElement(pPr, f'{ns}lnSpc')
                spcPct = etree.SubElement(lnSpc, f'{ns}spcPct')
                spcPct.set('val', '90000')
                buNone = etree.SubElement(pPr, f'{ns}buNone')

                r = etree.SubElement(p, f'{ns}r')
                rPr = etree.SubElement(r, f'{ns}rPr')
                rPr.set('lang', lang)
                rPr.set('sz', '2400')
                if bold:
                    rPr.set('b', '1')
                fill = etree.SubElement(rPr, f'{ns}solidFill')
                clr = etree.SubElement(fill, f'{ns}schemeClr')
                clr.set('val', 'bg1')
                latin = etree.SubElement(rPr, f'{ns}latin')
                latin.set('typeface', 'Montserrat')
                ea = etree.SubElement(rPr, f'{ns}ea')
                ea.set('typeface', 'Montserrat')
                cs = etree.SubElement(rPr, f'{ns}cs')
                cs.set('typeface', 'Montserrat')
                t_elem = etree.SubElement(r, f'{ns}t')
                t_elem.text = text

            add_author_line(txBody, "Выполнил: студент группы ИС-241", bold=False)
            add_author_line(txBody, "Бесчетнов Егор Павлович")
            add_author_line(txBody, "")
            add_author_line(txBody, "Руководитель: зав. каф. ПМиК, д.т.н., доцент", bold=False)
            add_author_line(txBody, "Нечта Иван Васильевич")

        elif full_text.strip() == "Выпускная квалификационная работа бакалавра":
            pass  # Keep as is


# ---- SLIDE 2: Цель и актуальность ----
slide2 = prs.slides[1]
set_title(slide2, "Цель и актуальность работы")
set_body_text(slide2, [
    "Цель: разработка программного приложения, реализующего полный цикл генерации клиентского интерфейса: от текстового запроса через локальную LLM и нормализацию до готового HTML-прототипа",
    "",
    "Актуальность:",
    "Разработчики без дизайнерских навыков испытывают трудности с UI-прототипированием",
    "Существующие ИИ-инструменты (Uizard, Galileo AI, Framer AI) требуют облака и платной подписки",
    "Зрелость экосистемы локальных LLM (2024-2025) делает офлайн-решение технически доступным",
    "Потребность в конфиденциальности данных в организациях",
])
set_page_number(slide2, 2)


# ---- SLIDE 3: Постановка задачи ----
slide3 = prs.slides[2]
set_title(slide3, "Задачи работы")
set_body_text(slide3, [
    "Проанализировать существующие инструменты и подходы к автоматизированному прототипированию UI",
    "Исследовать применимость локальных моделей llama3.2:3b, llama3.2:8b, gemma2:2b к генерации JSON",
    "Спроектировать формат JSON-спецификации интерфейса",
    "Разработать модуль нормализации ответов LLM, устойчивый к вариациям и ошибкам генерации",
    "Реализовать HTML-генератор с 4 типами компоновки (vertical, horizontal, grid, sidebar)",
    "Разработать графический интерфейс на CustomTkinter с тремя режимами генерации",
    "Реализовать режим анализа визуальных референсов (multimodal LLM)",
    "Провести функциональное тестирование и оценку качества",
])
set_page_number(slide3, 3)


# ---- SLIDE 4: Архитектура системы (was "Описание разработки") ----
slide4 = prs.slides[3]
set_title(slide4, "Архитектура системы")
set_body_text(slide4, [
    "Паттерн «LLM → JSON → HTML»: языковая модель формирует JSON-спецификацию, детерминированный генератор строит HTML",
    "5 модулей: GUI (main.py), Парсер (parser.py), Нормализатор, Генератор (generator.py), Логгер (logger.py)",
    "Контракт данных: {title, items, style, screens, extra_components}",
    "Graceful degradation: fallback при ошибках LLM",
])
set_page_number(slide4, 4)
# Add architecture diagram - below text
arch_img = os.path.join(IMG_DIR, "fig2_1_architecture_crop.png")
if os.path.exists(arch_img):
    add_image_to_slide(slide4, arch_img,
                       Emu(600000), Emu(4200000),
                       width=Emu(11000000))


# ---- SLIDE 5: Заключение -> repurpose for JSON spec ----
slide5 = prs.slides[4]
set_title(slide5, "JSON-спецификация интерфейса")
set_body_text(slide5, [
    "title — заголовок интерфейса (обязательное поле)",
    "items — массив строк, пункты главного меню (3-8 пунктов)",
    "style — визуальные параметры: bg_color, button_bg, text_color, layout",
    "screens — словарь вложенных экранов с title, content, components",
    "extra_components — интерактивные элементы: slider, checkbox, select, input",
    "layout: vertical | horizontal | grid | sidebar",
])
set_page_number(slide5, 5)
json_img = os.path.join(IMG_DIR, "fig2_2_json_spec_crop.png")
if os.path.exists(json_img):
    add_image_to_slide(slide5, json_img,
                       Emu(1500000), Emu(4200000),
                       width=Emu(9000000))


# Now we need more slides. We'll add them after slide 5 (before "Спасибо").
# The "Спасибо за внимание" slide is at index 5.

# We'll duplicate slide 1 (content slide) for additional slides
# Note: after adding slides, indices shift. Let's add all slides first.

# ---- Additional slides to add (will be appended after existing slides) ----

# We need to carefully manage the order. python-pptx adds slides at the end.
# At the end we'll reorder the sldIdLst.

# SLIDE 6: Нормализация
s6 = clone_slide(prs, 1)
set_title(s6, "Алгоритм нормализации")
set_body_text(s6, [
    "Зафиксировано более 40 вариантов имён полей от LLM (navigation, menu_items, buttons...)",
    "Шаг 1: Развёртывание обёртки ({\"interface\": {...}} → {...})",
    "Шаг 2: Поиск полей по словарям псевдонимов (18 вариантов для items)",
    "Шаг 3: Нормализация типов (строка → массив, список словарей → список строк)",
    "Шаг 4: Определение layout по ключевым словам текста запроса",
    "Шаг 5: Автогенерация screens и components при отсутствии",
    "Все функции идемпотентны — повторный вызов не меняет результат",
])
set_page_number(s6, 6)
resize_body(s6, new_width=Emu(7200000))
norm_img = os.path.join(IMG_DIR, "fig2_4_normalization_crop.png")
if os.path.exists(norm_img):
    add_image_to_slide(s6, norm_img,
                       Emu(7800000), Emu(1200000),
                       height=Emu(5400000))


# SLIDE 7: HTML-генератор
s7 = clone_slide(prs, 1)
set_title(s7, "HTML-генератор и режимы работы")
set_body_text(s7, [
    "HTMLGenerator: _build_css() + _build_html() + _build_js()",
    "4 типа компоновки: vertical (flexbox), horizontal (row), grid (CSS Grid), sidebar (fixed nav)",
    "Интерактивные компоненты: slider, checkbox, select, input",
    "Самодостаточный HTML без внешних зависимостей (8-25 КБ)",
    "",
    "3 режима генерации:",
    "Текстовый — описание на естественном языке → JSON → HTML",
    "Шаблонный — выбор из 6 готовых шаблонов с адаптацией через LLM",
    "Vision — анализ 1-3 скриншотов через llava:7b → JSON → HTML",
])
set_page_number(s7, 7)


# SLIDE 8: GUI
s8 = clone_slide(prs, 1)
set_title(s8, "Графический интерфейс приложения")
set_body_text(s8, [
    "CustomTkinter с тёмной темой, 3 зоны: ввод, управление, результат",
    "Двухпанельный макет: JSON-редактор + HTML-превью (tkinterweb)",
    "Панель шаблонов, история сессий, режим визуальных референсов",
])
set_page_number(s8, 8)
gui_img = os.path.join(IMG_DIR, "fig2_6_gui_crop.png")
if os.path.exists(gui_img):
    add_image_to_slide(s8, gui_img,
                       Emu(600000), Emu(3400000),
                       width=Emu(11000000))


# SLIDE 9: Результаты тестирования
s9 = clone_slide(prs, 1)
set_title(s9, "Результаты тестирования")
set_body_text(s9, [
    "12 тест-кейсов различной сложности, 3 попытки каждый",
    "Режим текста: 75% валидный JSON без нормализатора, 100% с ним",
    "Режим шаблонов: 5 из 6 адаптаций верны, fallback на базовый шаблон",
    "Режим Vision: 3 из 4 — JSON от llava:7b, 1 — fallback на текст",
    "Среднее время: 7.2 с (простые) — 14.8 с (сложные), макс. 28 с",
])
set_page_number(s9, 9)
resize_body(s9, new_width=Emu(5500000))
test_img = os.path.join(IMG_DIR, "fig4_1_test_results_crop.png")
if os.path.exists(test_img):
    add_image_to_slide(s9, test_img,
                       Emu(6000000), Emu(1300000),
                       width=Emu(5800000))


# SLIDE 10: Оценка качества
s10 = clone_slide(prs, 1)
set_title(s10, "Оценка качества интерфейсов")
set_body_text(s10, [
    "Корректность JSON: 4.2/5 (10 из 12 без нормализации)",
    "Содержательное соответствие: 4.1/5 (9 из 12 точно по запросу)",
    "Правильность layout: 4.3/5 (11 из 12 верно)",
    "Визуальное качество: 3.8/5 (склонность к тёмным схемам)",
    "Средний итог: 4.1 из 5 баллов",
    "",
    "llama3.2:8b: 4.4/5 при 22-38 с (прирост 0.3 балла при 3x времени)",
    "Fallback обеспечивает ненулевой результат в 100% случаев",
])
set_page_number(s10, 10)
resize_body(s10, new_width=Emu(6500000))
layouts_img = os.path.join(IMG_DIR, "fig4_3_layouts_crop.png")
if os.path.exists(layouts_img):
    add_image_to_slide(s10, layouts_img,
                       Emu(7000000), Emu(1200000),
                       height=Emu(5200000))


# SLIDE 11: Результат работы
s11 = clone_slide(prs, 1)
set_title(s11, "Демонстрация результата")
set_body_text(s11, [
    "Пример: аналитический дашборд с sidebar-компоновкой",
    "Генерация за ~14 с на CPU (llama3.2:3b, Q4_K_M, 2 ГБ RAM)",
])
set_page_number(s11, 11)
# Result screenshot
result_img = os.path.join(IMG_DIR, "fig4_5_result_crop.png")
if os.path.exists(result_img):
    add_image_to_slide(s11, result_img,
                       Emu(300000), Emu(2800000),
                       width=Emu(5800000))
comp_img = os.path.join(IMG_DIR, "fig4_4_components_crop.png")
if os.path.exists(comp_img):
    add_image_to_slide(s11, comp_img,
                       Emu(6300000), Emu(2800000),
                       width=Emu(5600000))


# SLIDE 12: Заключение
s12 = clone_slide(prs, 1)
set_title(s12, "Заключение")
set_body_text(s12, [
    "Все поставленные задачи решены в полном объёме",
    "Ключевой вклад — архитектурный паттерн «LLM → JSON → HTML»",
    "Модуль нормализации снизил долю непригодных ответов с ~40% до 5%",
    "Три режима генерации: текст, шаблоны, визуальные референсы",
    "Средний балл качества: 4.1/5, время генерации: 7-28 с (офлайн)",
    "Подход «малая модель + умный постпроцессор» эффективнее замены на крупную модель",
    "",
    "Направления развития: JSON Schema валидация, генерация React-компонентов, интеграция с Figma API, адаптивный дизайн",
])
set_page_number(s12, 12)


# ---- Reorder slides ----
# Current order after additions:
# 0: Title, 1: Цель, 2: Задачи, 3: Архитектура, 4: JSON-спец, 5: Спасибо
# 6: Нормализация, 7: HTML-генератор, 8: GUI, 9: Тестирование, 10: Оценка, 11: Результат, 12: Заключение

# Desired order:
# 0: Title, 1: Цель, 2: Задачи, 3: Архитектура, 4: JSON-спец,
# 5(6): Нормализация, 6(7): HTML-генератор, 7(8): GUI,
# 8(9): Тестирование, 9(10): Оценка, 10(11): Результат, 11(12): Заключение, 12(5): Спасибо

# Reorder sldIdLst
sldIdLst = prs.slides._sldIdLst
sldIds = list(sldIdLst)

# Current indices: 0,1,2,3,4,5,6,7,8,9,10,11,12
# Desired:         0,1,2,3,4,6,7,8,9,10,11,12,5
desired_order = [0, 1, 2, 3, 4, 6, 7, 8, 9, 10, 11, 12, 5]

reordered = [sldIds[i] for i in desired_order]

# Clear and re-add
for sid in list(sldIdLst):
    sldIdLst.remove(sid)
for sid in reordered:
    sldIdLst.append(sid)

prs.save(OUTPUT)
print(f"Saved to {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
